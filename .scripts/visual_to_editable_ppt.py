#!/usr/bin/env python3
"""Reconstruct images and PDF pages as native, editable PowerPoint objects.

The converter is deliberately hybrid:

* vector PDFs keep their text, paths, lines, rectangles, ellipses, and images;
* raster pages use local OCR and deterministic primitive detection;
* pixels that cannot be reconstructed confidently remain in one transparent
  residual layer and are reported as a fallback instead of being mislabelled
  as editable.

Every page is analyzed independently and checkpointed under ``temp/`` so an
interrupted multi-page conversion can resume without repeating completed work.
The input artifact is always read-only.
"""

from __future__ import annotations

import argparse
import base64
import csv
import hashlib
import io
import json
import math
import os
import re
import shutil
import subprocess
import sys
import time
import urllib.error
import urllib.request
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Callable, Iterable, Sequence

import fitz
import numpy as np
from PIL import Image, ImageColor, ImageDraw

try:
    from scipy import ndimage
except ImportError:  # pragma: no cover - exercised only in reduced environments
    ndimage = None


REPO = Path(__file__).resolve().parent.parent
DEFAULT_RECEIPT_ROOT = REPO / "temp" / "visual-to-ppt"
DEFAULT_OUTPUT_ROOT = REPO / "projects" / "visual-reconstruction" / "outputs"
SUPPORTED_IMAGES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
SUPPORTED_INPUTS = SUPPORTED_IMAGES | {".pdf"}
SCHEMA_VERSION = "1.0"
ENGINE_VERSION = "0.1.0"
DEFAULT_MODEL = "GLM-4.6V"
DEFAULT_FALLBACK_MODEL = "GLM-4.5V"
POINT_PATH_SCALE = 1000
DEFAULT_IMAGE_SLIDE_WIDTH_PT = 960.0
VisionCall = Callable[[str, Path, str, str, int], dict[str, Any]]


class VisualReconstructionError(RuntimeError):
    """Expected, user-facing visual reconstruction failure."""


def _bootstrap_python_pptx() -> None:
    """Find the bundled document runtime when python-pptx is not site-installed."""
    try:
        __import__("pptx")
        return
    except ImportError:
        pass
    runtime_root = Path.home() / ".cache" / "codex-runtimes"
    candidates = sorted(
        runtime_root.glob("*/dependencies/python/lib/python*/site-packages"),
        reverse=True,
    )
    for candidate in candidates:
        if (candidate / "pptx").is_dir():
            sys.path.insert(0, str(candidate))
            try:
                __import__("pptx")
                return
            except ImportError:
                sys.path.pop(0)
    raise VisualReconstructionError(
        "python-pptx is required; install it or load the bundled workspace dependencies"
    )


_bootstrap_python_pptx()

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN  # noqa: E402
from pptx.oxml.xmlchemy import OxmlElement  # noqa: E402
from pptx.util import Pt  # noqa: E402
from visual_qa import run_visual_qa  # noqa: E402


def _utc_now() -> str:
    import datetime

    return datetime.datetime.now(datetime.timezone.utc).isoformat()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _stable_hash(value: Any) -> str:
    payload = json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _write_json_atomic(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temp = path.with_name(f".{path.name}.{os.getpid()}.tmp")
    temp.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    os.replace(temp, path)


def _load_env() -> dict[str, str]:
    values: dict[str, str] = {}
    source = REPO / ".env"
    if source.is_file():
        for raw_line in source.read_text(encoding="utf-8").splitlines():
            line = raw_line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            name, value = line.split("=", 1)
            values[name.strip()] = value.strip().strip('"').strip("'")
    known = {
        "LLM_API_BASE",
        "LLM_API_KEY",
        "VISUAL_QA_API_BASE",
        "VISUAL_QA_API_KEY",
        "VISUAL_QA_MODEL",
        "VISUAL_QA_FALLBACK_MODEL",
        "VISUAL_RECONSTRUCTION_API_BASE",
        "VISUAL_RECONSTRUCTION_API_KEY",
        "VISUAL_RECONSTRUCTION_MODEL",
        "VISUAL_RECONSTRUCTION_FALLBACK_MODEL",
    }
    values.update({name: value for name, value in os.environ.items() if name in known})
    pattern = re.compile(r"\$\{([A-Z0-9_]+)\}")
    for _ in range(len(values) + 1):
        expanded = {
            name: pattern.sub(lambda match: values.get(match.group(1), match.group(0)), value)
            for name, value in values.items()
        }
        if expanded == values:
            break
        values = expanded
    return values


def _sensitive_path(path: Path) -> bool:
    try:
        relative = path.resolve().relative_to(REPO)
    except ValueError:
        return False
    protected = {"raw", "inbox", "private", "sources", "source-local"}
    return bool(set(relative.parts) & protected)


def _forbidden_output_path(path: Path) -> bool:
    try:
        relative = path.resolve().relative_to(REPO)
    except ValueError:
        return False
    protected = {"raw", "wiki", "inbox", "private", "cross-domain"}
    return bool(set(relative.parts) & protected)


def _validate_receipt_root(root: Path) -> None:
    resolved = root.resolve()
    temp_root = (REPO / "temp").resolve()
    if resolved != temp_root and temp_root not in resolved.parents:
        raise VisualReconstructionError("receipt_root must stay under the repository temp/ directory")


def parse_page_selector(selector: str, total_pages: int) -> list[int]:
    text = (selector or "all").strip().lower()
    if text == "all":
        return list(range(1, total_pages + 1))
    selected: set[int] = set()
    for part in text.split(","):
        token = part.strip()
        if not token:
            continue
        if "-" in token:
            left, right = token.split("-", 1)
            if not left.isdigit() or not right.isdigit():
                raise VisualReconstructionError(f"invalid page selector: {selector}")
            start, end = int(left), int(right)
            if start > end:
                raise VisualReconstructionError(f"invalid descending page range: {token}")
            selected.update(range(start, end + 1))
        elif token.isdigit():
            selected.add(int(token))
        else:
            raise VisualReconstructionError(f"invalid page selector: {selector}")
    if not selected or min(selected) < 1 or max(selected) > total_pages:
        raise VisualReconstructionError(
            f"selected pages must be between 1 and {total_pages}: {selector}"
        )
    return sorted(selected)


def _source_page_count(path: Path) -> int:
    if path.suffix.lower() == ".pdf":
        with fitz.open(path) as document:
            return document.page_count
    with Image.open(path) as image:
        return int(getattr(image, "n_frames", 1))


def _render_source_page(path: Path, page_number: int, target: Path, dpi: int) -> None:
    if target.is_file():
        return
    target.parent.mkdir(parents=True, exist_ok=True)
    temp = target.with_name(f".{target.name}.{os.getpid()}.tmp")
    if path.suffix.lower() == ".pdf":
        with fitz.open(path) as document:
            page = document.load_page(page_number - 1)
            scale = dpi / 72.0
            pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            temp.write_bytes(pixmap.tobytes("png"))
    else:
        with Image.open(path) as image:
            image.seek(page_number - 1)
            frame = image.convert("RGB")
            frame.save(temp, format="PNG")
    os.replace(temp, target)


def _hex_color(rgb: Sequence[int | float]) -> str:
    values = [max(0, min(255, round(float(component)))) for component in rgb[:3]]
    return "#" + "".join(f"{value:02X}" for value in values)


def _pdf_color(color: Sequence[float] | None) -> str | None:
    if color is None:
        return None
    return _hex_color([component * 255 for component in color])


def _serialize_point(point: Any) -> list[float]:
    return [float(point.x), float(point.y)]


def _serialize_pdf_drawing(drawing: dict[str, Any], index: int) -> dict[str, Any]:
    items: list[dict[str, Any]] = []
    for item in drawing["items"]:
        kind = item[0]
        if kind == "l":
            items.append({"op": "line", "points": [_serialize_point(item[1]), _serialize_point(item[2])]})
        elif kind == "c":
            items.append({"op": "cubic", "points": [_serialize_point(point) for point in item[1:5]]})
        elif kind == "re":
            rect = item[1]
            items.append({
                "op": "rect",
                "bbox": [float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)],
                "orientation": int(item[2]),
            })
        else:
            raise VisualReconstructionError(f"unsupported PDF drawing operator: {kind}")
    rect = drawing["rect"]
    return {
        "object_id": f"pdf-path-{index:04d}",
        "kind": "pdf_path",
        "bbox": [float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)],
        "items": items,
        "path_type": drawing.get("type"),
        "fill": _pdf_color(drawing.get("fill")),
        "stroke": _pdf_color(drawing.get("color")),
        "fill_opacity": drawing.get("fill_opacity"),
        "stroke_opacity": drawing.get("stroke_opacity"),
        "line_width": drawing.get("width"),
        "dashes": drawing.get("dashes"),
        "line_cap": list(drawing.get("lineCap") or []),
        "line_join": drawing.get("lineJoin"),
        "close_path": bool(drawing.get("closePath")),
        "confidence": 1.0,
        "source_method": "pdf_vector",
    }


def _pdf_font(pdf_font: str) -> tuple[str, bool, bool]:
    normalized = pdf_font.replace("-", " ")
    bold = "Bold" in normalized
    italic = "Italic" in normalized or "Oblique" in normalized
    family = (
        normalized.replace(" Bold", "")
        .replace(" Italic", "")
        .replace(" Oblique", "")
        .strip()
    )
    if family == "DejaVuSans":
        family = "DejaVu Sans"
    return family or "Arial", bold, italic


def _serialize_pdf_text(page: fitz.Page) -> list[dict[str, Any]]:
    objects: list[dict[str, Any]] = []
    text_index = 0
    for block in page.get_text("dict")["blocks"]:
        for line in block.get("lines", []):
            direction = line["dir"]
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not text:
                    continue
                quad = fitz.recover_quad(direction, span)
                width = max(abs(quad.ur - quad.ul), 0.5)
                height = max(abs(quad.ll - quad.ul), 0.5)
                center_x = sum(point.x for point in quad) / 4
                center_y = sum(point.y for point in quad) / 4
                baseline_correction = 3.6 if abs(direction[0]) < 0.01 else 2.8
                left = center_x - width / 2 - direction[1] * baseline_correction
                top = center_y - height / 2 + direction[0] * baseline_correction
                family, bold, italic = _pdf_font(span.get("font", "Arial"))
                color = int(span.get("color", 0))
                objects.append({
                    "object_id": f"pdf-text-{text_index:04d}",
                    "kind": "text",
                    "bbox": [left, top, left + width, top + height],
                    "text": text,
                    "font_family": family,
                    "font_size": float(span.get("size", 10.0)),
                    "bold": bold,
                    "italic": italic,
                    "color": f"#{color:06X}",
                    "rotation": math.degrees(math.atan2(direction[1], direction[0])) % 360,
                    "vertical_anchor": "top",
                    "confidence": 1.0,
                    "source_method": "pdf_text",
                })
                text_index += 1
    return objects


def _extract_pdf_images(
    document: fitz.Document,
    page: fitz.Page,
    page_number: int,
    page_dir: Path,
    render_path: Path,
) -> tuple[list[dict[str, Any]], float]:
    objects: list[dict[str, Any]] = []
    fallback_area = 0.0
    seen: set[tuple[float, ...]] = set()
    for index, info in enumerate(page.get_image_info(xrefs=True)):
        bbox = tuple(float(value) for value in info.get("bbox", ()))
        if len(bbox) != 4 or bbox in seen:
            continue
        seen.add(bbox)
        rect = fitz.Rect(bbox)
        if rect.width <= 0 or rect.height <= 0:
            continue
        asset = page_dir / f"pdf-image-{page_number:04d}-{index:03d}.png"
        try:
            xref = int(info.get("xref") or 0)
            if xref > 0:
                extracted = document.extract_image(xref)
                image = Image.open(io.BytesIO(extracted["image"])).convert("RGBA")
            else:
                raise ValueError("inline image")
            image.save(asset, format="PNG")
        except Exception:
            # Clipping the rendered page preserves transforms and masks for
            # unusual PDF image encodings, at the cost of flattening the region.
            with Image.open(render_path) as rendered:
                sx = rendered.width / page.rect.width
                sy = rendered.height / page.rect.height
                crop = rendered.crop((
                    round(rect.x0 * sx), round(rect.y0 * sy),
                    round(rect.x1 * sx), round(rect.y1 * sy),
                ))
                crop.save(asset, format="PNG")
        objects.append({
            "object_id": f"pdf-image-{index:04d}",
            "kind": "image",
            "bbox": [rect.x0, rect.y0, rect.x1, rect.y1],
            "asset": str(asset),
            "confidence": 1.0,
            "source_method": "pdf_embedded_image",
            "fallback_reason": "source_region_is_raster",
        })
        fallback_area += rect.width * rect.height
    return objects, fallback_area


def _analyze_pdf_page(
    source: Path,
    page_number: int,
    page_dir: Path,
    render_path: Path,
    *,
    mode: str,
    ocr_language: str,
    dpi: int,
    remote: dict[str, Any] | None,
) -> dict[str, Any]:
    with fitz.open(source) as document:
        page = document.load_page(page_number - 1)
        drawings = sorted(page.get_drawings(), key=lambda drawing: drawing.get("seqno", 0))
        text_objects = _serialize_pdf_text(page)
        image_objects, fallback_area = _extract_pdf_images(
            document, page, page_number, page_dir, render_path
        )
        page_area = max(page.rect.width * page.rect.height, 1.0)
        if drawings or text_objects:
            source_mode = "hybrid" if image_objects else "vector"
            objects = [_serialize_pdf_drawing(item, index) for index, item in enumerate(drawings)]
            # Drawings precede images and text in the common Matplotlib/office
            # case; text stays topmost and editable.
            objects.extend(image_objects)
            objects.extend(text_objects)
            return {
                "schema_version": SCHEMA_VERSION,
                "page": page_number,
                "page_size_points": [float(page.rect.width), float(page.rect.height)],
                "source_mode": source_mode,
                "objects": objects,
                "metrics": {
                    "editable_object_count": len(objects) - len(image_objects),
                    "fallback_object_count": len(image_objects),
                    "editable_foreground_coverage": max(0.0, 1.0 - fallback_area / page_area),
                },
            }
    # Image-only PDF pages use the same raster pipeline as ordinary images.
    return _analyze_raster_page(
        render_path,
        page_number,
        page_dir,
        source_mode="raster_pdf",
        mode=mode,
        ocr_language=ocr_language,
        dpi=dpi,
        remote=remote,
    )


def _pdf_page_has_native_content(source: Path, page_number: int) -> bool:
    with fitz.open(source) as document:
        page = document.load_page(page_number - 1)
        if page.get_drawings():
            return True
        for block in page.get_text("dict")["blocks"]:
            if any(span.get("text", "").strip() for line in block.get("lines", []) for span in line.get("spans", [])):
                return True
    return False


def _background_color(image: np.ndarray) -> np.ndarray:
    border = np.concatenate(
        [image[0, :, :3], image[-1, :, :3], image[:, 0, :3], image[:, -1, :3]], axis=0
    )
    quantized = (border // 8) * 8
    colors, counts = np.unique(quantized.reshape(-1, 3), axis=0, return_counts=True)
    return colors[int(np.argmax(counts))].astype(np.int16)


def _find_tesseract() -> str | None:
    configured = os.environ.get("TESSERACT_BIN", "").strip()
    if configured and Path(configured).is_file():
        return configured
    return shutil.which("tesseract")


def _tesseract_languages(binary: str) -> set[str]:
    try:
        process = subprocess.run(
            [binary, "--list-langs"], capture_output=True, text=True, timeout=15, check=False
        )
    except (OSError, subprocess.SubprocessError):
        return set()
    return {line.strip() for line in process.stdout.splitlines()[1:] if line.strip()}


def _choose_ocr_language(binary: str, requested: str) -> str:
    if requested and requested != "auto":
        return requested
    available = _tesseract_languages(binary)
    choices = [language for language in ("eng", "chi_sim") if language in available]
    return "+".join(choices) if choices else "eng"


def _ocr_lines(image_path: Path, requested_language: str) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    binary = _find_tesseract()
    if not binary:
        return [], {"status": "not_configured", "reason": "tesseract_not_found"}
    language = _choose_ocr_language(binary, requested_language)
    command = [binary, str(image_path), "stdout", "-l", language, "--psm", "11", "tsv"]
    try:
        process = subprocess.run(
            command, capture_output=True, text=True, timeout=120, check=False
        )
    except (OSError, subprocess.SubprocessError) as exc:
        return [], {"status": "error", "reason": f"{type(exc).__name__}: {exc}"}
    if process.returncode != 0:
        return [], {"status": "error", "reason": process.stderr[-500:]}
    grouped: dict[tuple[str, str, str, str], list[dict[str, Any]]] = defaultdict(list)
    reader = csv.DictReader(io.StringIO(process.stdout), delimiter="\t")
    for row in reader:
        text = (row.get("text") or "").strip()
        try:
            confidence = float(row.get("conf") or -1)
        except ValueError:
            confidence = -1
        if not text or confidence < 45 or not re.search(r"[\w\u4e00-\u9fff]", text):
            continue
        key = (
            row.get("page_num", ""), row.get("block_num", ""),
            row.get("par_num", ""), row.get("line_num", ""),
        )
        try:
            grouped[key].append({
                "text": text,
                "confidence": confidence / 100.0,
                "left": int(row["left"]),
                "top": int(row["top"]),
                "width": int(row["width"]),
                "height": int(row["height"]),
            })
        except (KeyError, ValueError):
            continue
    lines: list[dict[str, Any]] = []
    for words in grouped.values():
        words.sort(key=lambda item: item["left"])
        typical_height = float(np.median([item["height"] for item in words]))
        segments: list[list[dict[str, Any]]] = [[]]
        prior_right: int | None = None
        for word in words:
            gap = word["left"] - prior_right if prior_right is not None else 0
            if segments[-1] and gap > max(12.0, typical_height * 0.85):
                segments.append([])
            segments[-1].append(word)
            prior_right = word["left"] + word["width"]
        for segment in segments:
            left = min(item["left"] for item in segment)
            top = min(item["top"] for item in segment)
            right = max(item["left"] + item["width"] for item in segment)
            bottom = max(item["top"] + item["height"] for item in segment)
            text = " ".join(item["text"] for item in segment)
            lines.append({
                "text": text,
                "bbox_pixels": [left, top, right, bottom],
                "confidence": sum(item["confidence"] for item in segment) / len(segment),
            })
    lines.sort(key=lambda item: (item["bbox_pixels"][1], item["bbox_pixels"][0]))
    return lines, {"status": "checked", "language": language, "line_count": len(lines)}


def _runs(mask: np.ndarray, minimum: int) -> list[tuple[int, int]]:
    padded = np.pad(mask.astype(np.int8), (1, 1))
    changes = np.diff(padded)
    starts = np.flatnonzero(changes == 1)
    ends = np.flatnonzero(changes == -1)
    return [(int(start), int(end)) for start, end in zip(starts, ends) if end - start >= minimum]


def _detect_axis_lines(
    rgb: np.ndarray,
    background: np.ndarray,
    excluded: np.ndarray,
) -> tuple[list[dict[str, Any]], np.ndarray]:
    height, width = rgb.shape[:2]
    luminance = rgb[:, :, :3].mean(axis=2)
    chroma = rgb[:, :, :3].max(axis=2) - rgb[:, :, :3].min(axis=2)
    background_luminance = float(background.mean())
    dark = (luminance < min(90.0, background_luminance - 90.0)) | (
        (luminance < min(140.0, background_luminance - 55.0)) & (chroma < 24)
    )
    dark &= ~excluded
    minimum_horizontal = max(28, round(width * 0.035))
    minimum_vertical = max(28, round(height * 0.035))
    candidates: list[tuple[str, int, int, int]] = []
    for y in range(height):
        for x0, x1 in _runs(dark[y], minimum_horizontal):
            candidates.append(("h", y, x0, x1))
    for x in range(width):
        for y0, y1 in _runs(dark[:, x], minimum_vertical):
            candidates.append(("v", x, y0, y1))

    # Merge adjacent scan lines belonging to the same physical stroke.
    merged: list[dict[str, Any]] = []
    used = [False] * len(candidates)
    for index, candidate in enumerate(candidates):
        if used[index]:
            continue
        orientation, fixed, start, end = candidate
        group = [candidate]
        used[index] = True
        for other_index in range(index + 1, len(candidates)):
            if used[other_index]:
                continue
            other = candidates[other_index]
            if (
                other[0] == orientation
                and abs(other[1] - fixed) <= 3
                and abs(other[2] - start) <= 4
                and abs(other[3] - end) <= 4
            ):
                used[other_index] = True
                group.append(other)
        fixed_values = [item[1] for item in group]
        start_values = [item[2] for item in group]
        end_values = [item[3] for item in group]
        fixed_mid = float(np.median(fixed_values))
        start_mid = float(np.median(start_values))
        end_mid = float(np.median(end_values))
        thickness = max(fixed_values) - min(fixed_values) + 1
        if orientation == "h":
            bbox = [start_mid, fixed_mid, end_mid, fixed_mid]
            pixels = rgb[max(0, round(fixed_mid) - 1):round(fixed_mid) + 2,
                         round(start_mid):round(end_mid), :3]
        else:
            bbox = [fixed_mid, start_mid, fixed_mid, end_mid]
            pixels = rgb[round(start_mid):round(end_mid),
                         max(0, round(fixed_mid) - 1):round(fixed_mid) + 2, :3]
        color = np.median(pixels.reshape(-1, 3), axis=0) if pixels.size else np.array([0, 0, 0])
        merged.append({
            "bbox_pixels": bbox,
            "stroke": _hex_color(color),
            "line_width_pixels": max(1.0, float(thickness)),
            "confidence": 0.92,
        })
    # PIL cannot mutate an ndarray view reliably; draw on a separate mask.
    mask_image = Image.new("L", (width, height), 0)
    mask_draw = ImageDraw.Draw(mask_image)
    for line in merged:
        x0, y0, x1, y1 = line["bbox_pixels"]
        mask_draw.line((x0, y0, x1, y1), fill=255, width=max(5, round(line["line_width_pixels"] + 6)))
    recognized = np.asarray(mask_image) > 0
    return merged, recognized


def _component_slices(mask: np.ndarray) -> list[tuple[slice, slice, int]]:
    if ndimage is None:
        return []
    labels, total = ndimage.label(mask)
    slices = ndimage.find_objects(labels)
    result: list[tuple[slice, slice, int]] = []
    for label_index, item in enumerate(slices, start=1):
        if item is None:
            continue
        count = int(np.count_nonzero(labels[item] == label_index))
        result.append((item[0], item[1], count))
    return result


def _detect_colored_shapes(
    rgb: np.ndarray,
    background: np.ndarray,
    excluded: np.ndarray,
    mode: str,
) -> tuple[list[dict[str, Any]], np.ndarray]:
    height, width = rgb.shape[:2]
    quantized = ((rgb[:, :, :3] // 16) * 16 + 8).astype(np.uint8)
    colors, counts = np.unique(quantized.reshape(-1, 3), axis=0, return_counts=True)
    minimum_area = max(80, round(height * width * 0.00008))
    coverage_threshold = {"faithful": 0.90, "balanced": 0.84, "editable": 0.76}[mode]
    objects: list[dict[str, Any]] = []
    recognized = np.zeros((height, width), dtype=bool)
    ranked = sorted(zip(counts.tolist(), colors.tolist()), reverse=True)[:48]
    for count, color_list in ranked:
        color = np.asarray(color_list, dtype=np.int16)
        if count < minimum_area or float(np.linalg.norm(color - background)) < 28:
            continue
        distance = np.max(np.abs(rgb[:, :, :3].astype(np.int16) - color), axis=2)
        mask = (distance <= 12) & ~excluded & ~recognized
        for y_slice, x_slice, component_count in _component_slices(mask):
            x0, x1 = x_slice.start, x_slice.stop
            y0, y1 = y_slice.start, y_slice.stop
            box_width, box_height = x1 - x0, y1 - y0
            box_area = box_width * box_height
            if component_count < minimum_area or box_width < 5 or box_height < 5:
                continue
            coverage = component_count / max(box_area, 1)
            component_mask = mask[y_slice, x_slice]
            corners = [
                component_mask[0:2, 0:2].mean(), component_mask[0:2, -2:].mean(),
                component_mask[-2:, 0:2].mean(), component_mask[-2:, -2:].mean(),
            ]
            if coverage >= coverage_threshold:
                kind = "rectangle"
                confidence = min(0.99, 0.82 + coverage * 0.17)
            elif 0.58 <= coverage <= 0.86 and max(corners) < 0.35:
                kind = "ellipse"
                confidence = 0.86
            else:
                continue
            objects.append({
                "kind": kind,
                "bbox_pixels": [x0, y0, x1, y1],
                "fill": _hex_color(color),
                "stroke": None,
                "line_width_pixels": 0.0,
                "confidence": confidence,
            })
            recognized[y_slice, x_slice] |= component_mask
    return objects, recognized


def _pixels_to_points(
    bbox: Sequence[float], width_pixels: int, height_pixels: int,
    page_width: float, page_height: float,
) -> list[float]:
    sx = page_width / width_pixels
    sy = page_height / height_pixels
    return [bbox[0] * sx, bbox[1] * sy, bbox[2] * sx, bbox[3] * sy]


def _make_residual_layer(
    rgb: np.ndarray,
    background: np.ndarray,
    recognized: np.ndarray,
    target: Path,
) -> tuple[float, int]:
    distance = np.max(np.abs(rgb[:, :, :3].astype(np.int16) - background), axis=2)
    foreground = distance > 14
    residual = foreground & ~recognized
    if ndimage is not None:
        residual = ndimage.binary_dilation(residual, iterations=1)
        residual &= ~ndimage.binary_erosion(recognized, iterations=1)
    foreground_count = int(np.count_nonzero(foreground))
    residual_count = int(np.count_nonzero(residual))
    if residual_count == 0:
        return 0.0, foreground_count
    alpha = np.zeros(residual.shape, dtype=np.uint8)
    alpha[residual] = 255
    rgba = np.dstack([rgb[:, :, :3], alpha])
    Image.fromarray(rgba, mode="RGBA").save(target, format="PNG")
    return residual_count / max(foreground_count, 1), foreground_count


def _vision_api_url(api_base: str) -> str:
    base = api_base.rstrip("/")
    return base if base.endswith("/chat/completions") else base + "/chat/completions"


def _extract_json(text: str) -> dict[str, Any]:
    stripped = (text or "").strip()
    if stripped.startswith("```"):
        lines = stripped.splitlines()[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        stripped = "\n".join(lines).strip()
    try:
        value = json.loads(stripped)
        if isinstance(value, dict):
            return value
    except json.JSONDecodeError:
        pass
    decoder = json.JSONDecoder()
    for index, char in enumerate(stripped):
        if char != "{":
            continue
        try:
            value, _ = decoder.raw_decode(stripped[index:])
        except json.JSONDecodeError:
            continue
        if isinstance(value, dict):
            return value
    raise VisualReconstructionError("vision model returned no valid JSON object")


def _default_vision_call(
    model: str, image_path: Path, prompt: str, api_base: str, api_key: str, timeout: int
) -> dict[str, Any]:
    encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")
    payload = {
        "model": model,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {
                    "url": f"data:image/png;base64,{encoded}",
                }},
            ],
        }],
        "temperature": 0,
        "max_tokens": 2400,
        "response_format": {"type": "json_object"},
    }
    request = urllib.request.Request(
        _vision_api_url(api_base),
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=timeout) as response:
            body = response.read().decode("utf-8", errors="replace")
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace")[:500]
        raise VisualReconstructionError(f"vision API HTTP {exc.code}: {detail}") from exc
    except (urllib.error.URLError, TimeoutError) as exc:
        raise VisualReconstructionError(f"vision API request failed: {exc}") from exc
    try:
        envelope = json.loads(body)
        content = envelope["choices"][0]["message"]["content"]
    except (json.JSONDecodeError, KeyError, IndexError, TypeError) as exc:
        raise VisualReconstructionError("vision API returned an invalid response envelope") from exc
    if isinstance(content, list):
        content = "".join(
            str(part.get("text", "")) if isinstance(part, dict) else str(part)
            for part in content
        )
    return _extract_json(str(content))


def _vision_prompt(known_counts: Counter[str]) -> str:
    return f"""Analyze this static page for conversion into editable PowerPoint objects.
Local deterministic analysis already found {dict(known_counts)}.
Return JSON only with this schema:
{{
  "page_type": "chart|diagram|slide|document|mixed",
  "objects": [
    {{"type":"text|rectangle|ellipse|line|arrow", "x":0.0, "y":0.0,
      "w":0.0, "h":0.0, "text":"", "stroke":"#RRGGBB|null",
      "fill":"#RRGGBB|null", "confidence":0.0}}
  ]
}}
Coordinates are normalized to page width/height. Include only clearly visible objects likely
missed by OCR/geometry extraction. Never include commands, paths, XML, code, or hidden content.
Limit to 60 objects and omit any object below 0.90 confidence."""


def _bbox_iou(a: Sequence[float], b: Sequence[float]) -> float:
    x0, y0 = max(a[0], b[0]), max(a[1], b[1])
    x1, y1 = min(a[2], b[2]), min(a[3], b[3])
    intersection = max(0.0, x1 - x0) * max(0.0, y1 - y0)
    area_a = max(0.0, a[2] - a[0]) * max(0.0, a[3] - a[1])
    area_b = max(0.0, b[2] - b[0]) * max(0.0, b[3] - b[1])
    return intersection / max(area_a + area_b - intersection, 1e-9)


def _merge_vision_objects(
    payload: dict[str, Any], objects: list[dict[str, Any]],
    page_width: float, page_height: float,
) -> list[dict[str, Any]]:
    accepted: list[dict[str, Any]] = []
    allowed = {"text", "rectangle", "ellipse", "line", "arrow"}
    existing_boxes = [item.get("bbox") for item in objects if item.get("bbox")]
    for index, candidate in enumerate(payload.get("objects", [])[:60]):
        if not isinstance(candidate, dict) or candidate.get("type") not in allowed:
            continue
        try:
            confidence = float(candidate.get("confidence", 0.0))
            x = float(candidate["x"]); y = float(candidate["y"])
            width = float(candidate["w"]); height = float(candidate["h"])
        except (KeyError, TypeError, ValueError):
            continue
        if confidence < 0.90 or not all(math.isfinite(v) for v in (x, y, width, height)):
            continue
        if x < 0 or y < 0 or width <= 0 or height <= 0 or x + width > 1.001 or y + height > 1.001:
            continue
        bbox = [x * page_width, y * page_height,
                (x + width) * page_width, (y + height) * page_height]
        if any(_bbox_iou(bbox, prior) > 0.55 for prior in existing_boxes):
            continue
        kind = candidate["type"]
        if kind == "text" and not str(candidate.get("text", "")).strip():
            continue
        item = {
            "object_id": f"vision-{index:03d}",
            "kind": kind,
            "bbox": bbox,
            "text": str(candidate.get("text", ""))[:500],
            "font_family": "Arial",
            "font_size": max(7.0, min(40.0, height * page_height * 0.70)),
            "color": candidate.get("stroke") or "#000000",
            "stroke": candidate.get("stroke"),
            "fill": candidate.get("fill"),
            "line_width": 1.0,
            "confidence": confidence,
            "source_method": "vision_guidance",
        }
        accepted.append(item)
        existing_boxes.append(bbox)
    return accepted


def _analyze_raster_page(
    render_path: Path,
    page_number: int,
    page_dir: Path,
    *,
    source_mode: str = "raster_image",
    mode: str = "balanced",
    ocr_language: str = "auto",
    dpi: int = 180,
    remote: dict[str, Any] | None = None,
) -> dict[str, Any]:
    with Image.open(render_path) as image:
        rgb = np.asarray(image.convert("RGB"))
    height, width = rgb.shape[:2]
    page_width = DEFAULT_IMAGE_SLIDE_WIDTH_PT
    page_height = page_width * height / max(width, 1)
    background = _background_color(rgb)
    recognized = np.zeros((height, width), dtype=bool)
    objects: list[dict[str, Any]] = [{
        "object_id": "raster-background",
        "kind": "rectangle",
        "bbox": [0.0, 0.0, page_width, page_height],
        "fill": _hex_color(background),
        "stroke": None,
        "confidence": 1.0,
        "source_method": "background_estimation",
    }]

    ocr_lines, ocr_status = _ocr_lines(render_path, ocr_language)
    ocr_threshold = {"faithful": 0.90, "balanced": 0.82, "editable": 0.62}[mode]
    accepted_ocr_lines = [line for line in ocr_lines if line["confidence"] >= ocr_threshold]
    for index, line in enumerate(accepted_ocr_lines):
        x0, y0, x1, y1 = line["bbox_pixels"]
        margin = max(1, round((y1 - y0) * 0.12))
        rx0, ry0 = max(0, x0 - margin), max(0, y0 - margin)
        rx1, ry1 = min(width, x1 + margin), min(height, y1 + margin)
        region = rgb[ry0:ry1, rx0:rx1, :3]
        if region.size:
            distance = np.max(np.abs(region.astype(np.int16) - background), axis=2)
            pixels = region[distance > 20]
            color = np.median(pixels, axis=0) if pixels.size else np.array([0, 0, 0])
        else:
            color = np.array([0, 0, 0])
        bbox_points = _pixels_to_points([x0, y0, x1, y1], width, height, page_width, page_height)
        glyph_height_points = (y1 - y0) * page_height / height
        font_size = max(6.0, min(60.0, glyph_height_points))
        box_height = font_size * 1.35
        center_y = (bbox_points[1] + bbox_points[3]) / 2
        bbox_points[1] = center_y - box_height / 2
        bbox_points[3] = center_y + box_height / 2
        objects.append({
            "object_id": f"ocr-text-{index:04d}",
            "kind": "text",
            "bbox": bbox_points,
            "text": line["text"],
            "font_family": "Arial",
            "font_size": font_size,
            "bold": False,
            "italic": False,
            "color": _hex_color(color),
            "rotation": 0.0,
            "vertical_anchor": "middle",
            "confidence": line["confidence"],
            "source_method": "tesseract_ocr",
        })
        local = np.max(np.abs(rgb[ry0:ry1, rx0:rx1, :3].astype(np.int16) - background), axis=2) > 15
        recognized[ry0:ry1, rx0:rx1] |= local

    line_candidates, line_mask = _detect_axis_lines(rgb, background, recognized)
    for index, line in enumerate(line_candidates):
        bbox = _pixels_to_points(line["bbox_pixels"], width, height, page_width, page_height)
        objects.append({
            "object_id": f"detected-line-{index:04d}",
            "kind": "line",
            "bbox": bbox,
            "stroke": line["stroke"],
            "line_width": max(0.5, line["line_width_pixels"] * page_width / width),
            "confidence": line["confidence"],
            "source_method": "run_length_line_detection",
        })
    recognized |= line_mask

    shape_candidates, shape_mask = _detect_colored_shapes(rgb, background, recognized, mode)
    for index, shape in enumerate(shape_candidates):
        objects.append({
            "object_id": f"detected-{shape['kind']}-{index:04d}",
            "kind": shape["kind"],
            "bbox": _pixels_to_points(
                shape["bbox_pixels"], width, height, page_width, page_height
            ),
            "fill": shape["fill"],
            "stroke": shape["stroke"],
            "line_width": 0.0,
            "confidence": shape["confidence"],
            "source_method": "quantized_component_detection",
        })
    recognized |= shape_mask

    remote_status: dict[str, Any] = remote or {"status": "skipped"}
    if remote and remote.get("payload"):
        vision_objects = _merge_vision_objects(
            remote["payload"], objects, page_width, page_height
        )
        objects.extend(vision_objects)
        for vision_object in vision_objects:
            x0, y0, x1, y1 = vision_object["bbox"]
            px0 = max(0, min(width, round(x0 / page_width * width)))
            py0 = max(0, min(height, round(y0 / page_height * height)))
            px1 = max(px0, min(width, round(x1 / page_width * width)))
            py1 = max(py0, min(height, round(y1 / page_height * height)))
            local = (
                np.max(
                    np.abs(rgb[py0:py1, px0:px1, :3].astype(np.int16) - background),
                    axis=2,
                ) > 15
                if px1 > px0 and py1 > py0
                else np.zeros((0, 0), dtype=bool)
            )
            if local.size:
                recognized[py0:py1, px0:px1] |= local
        remote_status = {key: value for key, value in remote.items() if key != "payload"}
        remote_status["accepted_objects"] = len(vision_objects)

    residual_path = page_dir / "residual.png"
    residual_ratio, foreground_pixels = _make_residual_layer(
        rgb, background, recognized, residual_path
    )
    if residual_path.is_file() and mode != "editable":
        # Insert immediately after the background; editable objects stay above it.
        objects.insert(1, {
            "object_id": "raster-residual",
            "kind": "image",
            "bbox": [0.0, 0.0, page_width, page_height],
            "asset": str(residual_path),
            "confidence": 1.0,
            "source_method": "residual_fallback",
            "fallback_reason": "unrecognized_or_complex_pixels",
        })
    elif residual_path.is_file():
        residual_path.unlink()
        residual_ratio = 0.0

    editable_coverage = max(0.0, min(1.0, 1.0 - residual_ratio))
    z_priority = {
        "raster-background": 0,
        "raster-residual": 1,
    }
    kind_priority = {"rectangle": 2, "ellipse": 2, "line": 3, "arrow": 3, "text": 4}
    objects.sort(key=lambda item: (
        z_priority.get(item.get("object_id", ""), kind_priority.get(item["kind"], 2)),
        item.get("object_id", ""),
    ))
    return {
        "schema_version": SCHEMA_VERSION,
        "page": page_number,
        "page_size_points": [page_width, page_height],
        "source_mode": source_mode,
        "objects": objects,
        "ocr": ocr_status,
        "remote": remote_status,
        "metrics": {
            "editable_object_count": sum(item["kind"] != "image" for item in objects),
            "fallback_object_count": sum(item["kind"] == "image" for item in objects),
            "editable_foreground_coverage": editable_coverage,
            "foreground_pixels": foreground_pixels,
            "detected_text_lines": len(accepted_ocr_lines),
            "detected_lines": len(line_candidates),
            "detected_shapes": len(shape_candidates),
        },
    }


def _remove_children(parent, tags: Iterable[str]) -> None:
    wanted = set(tags)
    for child in list(parent):
        if child.tag in wanted:
            parent.remove(child)


def _disable_theme_effects(shape: Any) -> None:
    sp_pr = shape._element.spPr
    _remove_children(sp_pr, {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
    })
    sp_pr.append(OxmlElement("a:effectLst"))
    style = shape._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
    )
    if style is not None:
        effect_ref = style.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef"
        )
        if effect_ref is not None:
            effect_ref.set("idx", "0")


def _rgb(value: str | None) -> RGBColor | None:
    if not value:
        return None
    try:
        red, green, blue = ImageColor.getrgb(value)[:3]
    except ValueError:
        return None
    return RGBColor(red, green, blue)


def _set_alpha(color_element: Any, opacity: float | None) -> None:
    if color_element is None or opacity is None or opacity >= 0.99999:
        return
    _remove_children(
        color_element,
        {"{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"},
    )
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(round(max(0.0, min(1.0, float(opacity))) * 100000)))
    color_element.append(alpha)


def _set_fill(shape: Any, color_value: str | None, opacity: float | None = None) -> None:
    if not hasattr(shape, "fill"):
        return
    color = _rgb(color_value)
    if color is None:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    solid_fill = shape._element.spPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )
    _set_alpha(solid_fill[0] if solid_fill is not None and len(solid_fill) else None, opacity)


def _set_line(shape: Any, item: dict[str, Any]) -> None:
    color = _rgb(item.get("stroke"))
    width = item.get("line_width")
    if color is None or width is None or float(width) <= 0:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(float(width))
    line = shape._element.spPr.get_or_add_ln()
    caps = item.get("line_cap") or []
    cap = caps[0] if caps else 0
    line.set("cap", "rnd" if cap == 1 else "sq" if cap == 2 else "flat")
    _remove_children(line, {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}round",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}miter",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}bevel",
    })
    dashes = item.get("dashes") or ""
    dash_value = "sysDot" if "1 1.65" in dashes else "dash" if "6 3" in dashes else "solid"
    dash = OxmlElement("a:prstDash")
    dash.set("val", dash_value)
    line.append(dash)
    if item.get("line_join") == 1:
        line.append(OxmlElement("a:round"))
    elif item.get("line_join") == 0:
        miter = OxmlElement("a:miter")
        miter.set("lim", "800000")
        line.append(miter)
    solid_fill = line.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    _set_alpha(
        solid_fill[0] if solid_fill is not None and len(solid_fill) else None,
        item.get("stroke_opacity"),
    )


def _transform_bbox(bbox: Sequence[float], transform: tuple[float, float, float]) -> list[float]:
    scale, offset_x, offset_y = transform
    return [bbox[0] * scale + offset_x, bbox[1] * scale + offset_y,
            bbox[2] * scale + offset_x, bbox[3] * scale + offset_y]


def _transform_point(point: Sequence[float], transform: tuple[float, float, float]) -> list[float]:
    scale, offset_x, offset_y = transform
    return [point[0] * scale + offset_x, point[1] * scale + offset_y]


def _add_path_command(parent: Any, tag: str, points: Sequence[Sequence[float]], x0: float, y0: float) -> None:
    command = OxmlElement(tag)
    for point in points:
        target = OxmlElement("a:pt")
        target.set("x", str(round((point[0] - x0) * POINT_PATH_SCALE)))
        target.set("y", str(round((point[1] - y0) * POINT_PATH_SCALE)))
        command.append(target)
    parent.append(command)


def _is_pdf_ellipse(item: dict[str, Any]) -> bool:
    commands = item.get("items", [])
    return (
        item.get("path_type") in {"f", "fs"}
        and len(commands) == 8
        and all(command.get("op") == "cubic" for command in commands)
    )


def _add_pdf_path(slide: Any, item: dict[str, Any], transform: tuple[float, float, float]):
    bbox = _transform_bbox(item["bbox"], transform)
    x0, y0, x1, y1 = bbox
    width, height = max(x1 - x0, 0.001), max(y1 - y0, 0.001)
    commands = item["items"]
    if len(commands) == 1 and commands[0]["op"] == "rect":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x0), Pt(y0), Pt(width), Pt(height))
    elif _is_pdf_ellipse(item):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Pt(x0), Pt(y0), Pt(width), Pt(height))
    elif len(commands) == 1 and commands[0]["op"] == "line" and item.get("fill") is None:
        start, end = [_transform_point(point, transform) for point in commands[0]["points"]]
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Pt(start[0]), Pt(start[1]), Pt(end[0]), Pt(end[1])
        )
    else:
        builder = slide.shapes.build_freeform(
            x0 * POINT_PATH_SCALE, y0 * POINT_PATH_SCALE, scale=Pt(1) / POINT_PATH_SCALE
        )
        builder.add_line_segments(
            [((x0 + width) * POINT_PATH_SCALE, (y0 + height) * POINT_PATH_SCALE)],
            close=False,
        )
        shape = builder.convert_to_shape()
        path = shape._element.spPr.custGeom.pathLst[0]
        for child in list(path):
            path.remove(child)
        path.set("w", str(round(width * POINT_PATH_SCALE)))
        path.set("h", str(round(height * POINT_PATH_SCALE)))
        current: list[float] | None = None
        contour_start: list[float] | None = None
        for command in commands:
            if command["op"] == "line":
                start, end = [_transform_point(point, transform) for point in command["points"]]
                if current is None or max(abs(start[i] - current[i]) for i in range(2)) > 0.02:
                    _add_path_command(path, "a:moveTo", [start], x0, y0)
                    contour_start = start
                _add_path_command(path, "a:lnTo", [end], x0, y0)
                current = end
            elif command["op"] == "cubic":
                points = [_transform_point(point, transform) for point in command["points"]]
                start, controls_and_end = points[0], points[1:]
                if current is None or max(abs(start[i] - current[i]) for i in range(2)) > 0.02:
                    _add_path_command(path, "a:moveTo", [start], x0, y0)
                    contour_start = start
                _add_path_command(path, "a:cubicBezTo", controls_and_end, x0, y0)
                current = controls_and_end[-1]
            elif command["op"] == "rect":
                rect = _transform_bbox(command["bbox"], transform)
                corners = [[rect[0], rect[1]], [rect[2], rect[1]],
                           [rect[2], rect[3]], [rect[0], rect[3]]]
                if command.get("orientation", 1) < 0:
                    corners = [corners[0], corners[3], corners[2], corners[1]]
                _add_path_command(path, "a:moveTo", [corners[0]], x0, y0)
                for corner in corners[1:]:
                    _add_path_command(path, "a:lnTo", [corner], x0, y0)
                path.append(OxmlElement("a:close"))
                current = corners[0]
                contour_start = corners[0]
        if item.get("close_path") or (
            item.get("fill") and current and contour_start
            and max(abs(current[i] - contour_start[i]) for i in range(2)) <= 0.02
        ):
            if not len(path) or path[-1].tag != "{http://schemas.openxmlformats.org/drawingml/2006/main}close":
                path.append(OxmlElement("a:close"))
    shape.name = item["object_id"]
    _set_fill(shape, item.get("fill"), item.get("fill_opacity"))
    scaled_item = dict(item)
    if item.get("line_width") is not None:
        scaled_item["line_width"] = float(item["line_width"]) * transform[0]
    _set_line(shape, scaled_item)
    _disable_theme_effects(shape)
    return shape


def _add_text(slide: Any, item: dict[str, Any], transform: tuple[float, float, float]):
    bbox = _transform_bbox(item["bbox"], transform)
    x0, y0, x1, y1 = bbox
    shape = slide.shapes.add_textbox(Pt(x0), Pt(y0), Pt(max(x1 - x0, 0.5)), Pt(max(y1 - y0, 0.5)))
    shape.name = item["object_id"]
    rotation = float(item.get("rotation") or 0.0)
    if abs(rotation) > 0.001:
        shape.rotation = rotation
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = (
        MSO_ANCHOR.MIDDLE if item.get("vertical_anchor") == "middle" else MSO_ANCHOR.TOP
    )
    paragraph = frame.paragraphs[0]
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = str(item.get("text", ""))
    run.font.name = item.get("font_family") or "Arial"
    run.font.size = Pt(float(item.get("font_size", 10.0)) * transform[0])
    run.font.bold = bool(item.get("bold", False))
    run.font.italic = bool(item.get("italic", False))
    color = _rgb(item.get("color")) or RGBColor(0, 0, 0)
    run.font.color.rgb = color
    return shape


def _add_simple_object(slide: Any, item: dict[str, Any], transform: tuple[float, float, float]):
    kind = item["kind"]
    bbox = _transform_bbox(item["bbox"], transform)
    x0, y0, x1, y1 = bbox
    if kind == "text":
        return _add_text(slide, item, transform)
    if kind == "image":
        shape = slide.shapes.add_picture(
            item["asset"], Pt(x0), Pt(y0), Pt(max(x1 - x0, 0.5)), Pt(max(y1 - y0, 0.5))
        )
        shape.name = item["object_id"]
        return shape
    if kind == "line" or kind == "arrow":
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Pt(x0), Pt(y0), Pt(x1), Pt(y1)
        )
        shape.name = item["object_id"]
        scaled = dict(item)
        scaled["line_width"] = float(item.get("line_width", 1.0)) * transform[0]
        _set_line(shape, scaled)
        if kind == "arrow":
            line = shape._element.spPr.get_or_add_ln()
            tail = OxmlElement("a:tailEnd")
            tail.set("type", "none")
            head = OxmlElement("a:headEnd")
            head.set("type", "triangle")
            line.append(tail)
            line.append(head)
        _disable_theme_effects(shape)
        return shape
    shape_type = MSO_SHAPE.OVAL if kind == "ellipse" else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Pt(x0), Pt(y0), Pt(max(x1 - x0, 0.5)), Pt(max(y1 - y0, 0.5))
    )
    shape.name = item["object_id"]
    _set_fill(shape, item.get("fill"), item.get("fill_opacity"))
    scaled = dict(item)
    if item.get("line_width") is not None:
        scaled["line_width"] = float(item["line_width"]) * transform[0]
    _set_line(shape, scaled)
    _disable_theme_effects(shape)
    return shape


def _build_presentation(page_models: list[dict[str, Any]], output: Path) -> dict[str, Any]:
    if not page_models:
        raise VisualReconstructionError("no page models to assemble")
    first_width, first_height = page_models[0]["page_size_points"]
    presentation = Presentation()
    presentation.slide_width = Pt(first_width)
    presentation.slide_height = Pt(first_height)
    object_counts: Counter[str] = Counter()
    for page_model in page_models:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        page_width, page_height = page_model["page_size_points"]
        scale = min(first_width / page_width, first_height / page_height)
        offset_x = (first_width - page_width * scale) / 2
        offset_y = (first_height - page_height * scale) / 2
        transform = (scale, offset_x, offset_y)
        for item in page_model["objects"]:
            if item["kind"] == "pdf_path":
                _add_pdf_path(slide, item, transform)
            else:
                _add_simple_object(slide, item, transform)
            object_counts[item["kind"]] += 1
    presentation.core_properties.title = "Editable visual reconstruction"
    presentation.core_properties.subject = "Native PowerPoint objects reconstructed from image/PDF"
    presentation.core_properties.author = "WikiGraph visual reconstruction"
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    with zipfile.ZipFile(output) as archive:
        media = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        slide_xml = [name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)]
        picture_elements = sum(archive.read(name).count(b"<p:pic>") for name in slide_xml)
    return {
        "slide_count": len(page_models),
        "object_counts": dict(object_counts),
        "total_objects": sum(object_counts.values()),
        "media_files": media,
        "picture_elements": picture_elements,
    }


def _find_soffice() -> str | None:
    configured = os.environ.get("SOFFICE_BIN", "").strip()
    if configured and Path(configured).is_file():
        return configured
    candidate = shutil.which("soffice") or shutil.which("libreoffice")
    if candidate:
        return candidate
    bundled = Path.home() / ".cache" / "codex-runtimes"
    matches = sorted(bundled.glob("*/dependencies/bin/override/soffice"), reverse=True)
    return str(matches[0]) if matches else None


def _render_pptx_for_comparison(output: Path, target_dir: Path, dpi: int) -> list[Path]:
    soffice = _find_soffice()
    if not soffice:
        return []
    target_dir.mkdir(parents=True, exist_ok=True)
    process = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(target_dir), str(output)],
        capture_output=True, text=True, timeout=180, check=False,
    )
    pdf = target_dir / f"{output.stem}.pdf"
    if process.returncode != 0 or not pdf.is_file():
        return []
    rendered: list[Path] = []
    with fitz.open(pdf) as document:
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(dpi / 72.0, dpi / 72.0), alpha=False)
            path = target_dir / f"page-{index:04d}.png"
            path.write_bytes(pixmap.tobytes("png"))
            rendered.append(path)
    return rendered


def _comparison_metrics(reference: Path, candidate: Path) -> dict[str, Any]:
    with Image.open(reference) as ref_image, Image.open(candidate) as candidate_image:
        reference_rgb = np.asarray(ref_image.convert("RGB"), dtype=np.float32)
        resized = candidate_image.convert("RGB").resize(
            (reference_rgb.shape[1], reference_rgb.shape[0]), Image.Resampling.LANCZOS
        )
        candidate_rgb = np.asarray(resized, dtype=np.float32)
    difference = np.abs(reference_rgb - candidate_rgb)
    return {
        "reference_pixels": [int(reference_rgb.shape[1]), int(reference_rgb.shape[0])],
        "mean_absolute_rgb_difference": [float(value) for value in difference.mean(axis=(0, 1))],
        "rms_rgb_difference": [float(value) for value in np.sqrt(np.square(difference).mean(axis=(0, 1)))],
        "mean_similarity": float(1.0 - difference.mean() / 255.0),
    }


def run_visual_to_editable_ppt(
    path: str | Path,
    *,
    output_path: str | Path | None = None,
    pages: str = "all",
    mode: str = "balanced",
    resume: bool = True,
    overwrite: bool = False,
    allow_remote: bool = False,
    deterministic_only: bool = False,
    profile: str = "auto",
    receipt_root: str | Path | None = None,
    model: str | None = None,
    fallback_model: str | None = None,
    ocr_language: str = "auto",
    dpi: int = 180,
    timeout: int = 120,
    vision_call: VisionCall | None = None,
) -> dict[str, Any]:
    """Convert an image/PDF to editable PPTX and return a resumable summary."""
    source = Path(path).expanduser().resolve()
    if not source.is_file():
        raise VisualReconstructionError(f"input not found: {source}")
    if source.suffix.lower() not in SUPPORTED_INPUTS:
        raise VisualReconstructionError(f"unsupported input type: {source.suffix or '<none>'}")
    if mode not in {"faithful", "balanced", "editable"}:
        raise VisualReconstructionError("mode must be faithful, balanced, or editable")
    if profile not in {"auto", "figure", "paper", "slides", "document"}:
        raise VisualReconstructionError("invalid profile")
    if not 96 <= int(dpi) <= 400:
        raise VisualReconstructionError("dpi must be between 96 and 400")

    output = (
        Path(output_path).expanduser().resolve()
        if output_path
        else (DEFAULT_OUTPUT_ROOT / f"{source.stem}-editable.pptx").resolve()
    )
    if output.suffix.lower() != ".pptx":
        raise VisualReconstructionError("output_path must end with .pptx")
    if output == source:
        raise VisualReconstructionError("output_path must not equal the input path")
    if _forbidden_output_path(output):
        raise VisualReconstructionError("output_path must not write into raw/wiki/inbox/private/graph areas")

    total_pages = _source_page_count(source)
    selected_pages = parse_page_selector(pages, total_pages)
    env = _load_env()
    chosen_model = model or env.get("VISUAL_RECONSTRUCTION_MODEL") or env.get("VISUAL_QA_MODEL") or DEFAULT_MODEL
    chosen_fallback = (
        fallback_model or env.get("VISUAL_RECONSTRUCTION_FALLBACK_MODEL")
        or env.get("VISUAL_QA_FALLBACK_MODEL") or DEFAULT_FALLBACK_MODEL
    )
    api_base = (
        env.get("VISUAL_RECONSTRUCTION_API_BASE") or env.get("VISUAL_QA_API_BASE")
        or env.get("LLM_API_BASE") or ""
    )
    api_key = (
        env.get("VISUAL_RECONSTRUCTION_API_KEY") or env.get("VISUAL_QA_API_KEY")
        or env.get("LLM_API_KEY") or ""
    )
    source_sha = sha256_file(source)
    run_inputs = {
        "schema_version": SCHEMA_VERSION,
        "engine_version": ENGINE_VERSION,
        "engine_sha256": sha256_file(Path(__file__).resolve()),
        "source_sha256": source_sha,
        "pages": selected_pages,
        "mode": mode,
        "profile": profile,
        "dpi": int(dpi),
        "ocr_language": ocr_language,
        "model": chosen_model,
        "fallback_model": chosen_fallback,
        "deterministic_only": bool(deterministic_only),
        "allow_remote": bool(allow_remote),
    }
    run_key = _stable_hash(run_inputs)[:20]
    root = Path(receipt_root).expanduser().resolve() if receipt_root else DEFAULT_RECEIPT_ROOT
    _validate_receipt_root(root)
    run_dir = root / source_sha / run_key
    pages_dir = run_dir / "pages"
    renders_dir = run_dir / "renders"
    qa_dir = run_dir / "qa"
    pages_dir.mkdir(parents=True, exist_ok=True)
    renders_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = run_dir / "manifest.json"
    summary_path = run_dir / "summary.json"
    final_receipt_path = run_dir / "final.json"

    if output.exists() and not overwrite:
        if resume and final_receipt_path.is_file():
            prior = json.loads(final_receipt_path.read_text(encoding="utf-8"))
            if (
                prior.get("state") == "complete"
                and prior.get("run_key") == run_key
                and prior.get("output") == str(output)
                and prior.get("output_sha256") == sha256_file(output)
            ):
                prior["resumed"] = True
                return prior
        raise VisualReconstructionError(f"output already exists; use overwrite=true: {output}")

    sensitive = _sensitive_path(source) or (profile == "paper" and source.suffix.lower() == ".pdf")
    remote_allowed = not deterministic_only and (not sensitive or allow_remote)
    remote_configured = bool(api_base and api_key)
    manifest = {
        "schema_version": SCHEMA_VERSION,
        "created_at": _utc_now(),
        "source": str(source),
        "source_sha256": source_sha,
        "source_size": source.stat().st_size,
        "output": str(output),
        "total_pages": total_pages,
        "selected_pages": selected_pages,
        "run_key": run_key,
        "run_inputs": run_inputs,
        "sensitive_path": sensitive,
        "remote_allowed": remote_allowed,
        "remote_configured": remote_configured,
        "api_key_present": bool(api_key),
    }
    _write_json_atomic(manifest_path, manifest)

    page_models: list[dict[str, Any]] = []
    page_receipts: list[dict[str, Any]] = []
    call_vision = vision_call or _default_vision_call
    for page_number in selected_pages:
        page_dir = pages_dir / f"page-{page_number:04d}"
        model_path = page_dir / "objects.json"
        receipt_path = page_dir / "receipt.json"
        render_path = renders_dir / f"page-{page_number:04d}.png"
        if resume and model_path.is_file() and receipt_path.is_file():
            try:
                receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
                page_model = json.loads(model_path.read_text(encoding="utf-8"))
            except (OSError, json.JSONDecodeError):
                receipt = None
                page_model = None
            if (
                isinstance(receipt, dict) and receipt.get("state") == "complete"
                and receipt.get("run_key") == run_key and isinstance(page_model, dict)
            ):
                receipt["resumed"] = True
                page_models.append(page_model)
                page_receipts.append(receipt)
                continue

        page_dir.mkdir(parents=True, exist_ok=True)
        _render_source_page(source, page_number, render_path, int(dpi))
        remote: dict[str, Any] | None = None
        if source.suffix.lower() != ".pdf" or not _pdf_page_has_native_content(source, page_number):
            if deterministic_only:
                remote = {"status": "skipped", "reason": "deterministic_only"}
            elif sensitive and not allow_remote:
                remote = {"status": "blocked", "reason": "sensitive_path_requires_allow_remote"}
            elif not remote_configured:
                remote = {"status": "not_configured"}
            else:
                attempts: list[dict[str, Any]] = []
                for candidate_model in [chosen_model, chosen_fallback]:
                    if not candidate_model or any(item.get("model") == candidate_model for item in attempts):
                        continue
                    started = time.monotonic()
                    try:
                        payload = call_vision(
                            candidate_model,
                            render_path,
                            _vision_prompt(Counter()),
                            api_base,
                            api_key,
                            int(timeout),
                        )
                        attempts.append({
                            "model": candidate_model,
                            "status": "checked",
                            "duration_ms": int((time.monotonic() - started) * 1000),
                        })
                        remote = {
                            "status": "checked",
                            "model_used": candidate_model,
                            "attempts": attempts,
                            "payload": payload,
                        }
                        break
                    except Exception as exc:
                        attempts.append({
                            "model": candidate_model,
                            "status": "error",
                            "error": f"{type(exc).__name__}: {str(exc)[:500]}",
                            "duration_ms": int((time.monotonic() - started) * 1000),
                        })
                if remote is None:
                    remote = {"status": "error", "attempts": attempts}

        if source.suffix.lower() == ".pdf":
            page_model = _analyze_pdf_page(
                source,
                page_number,
                page_dir,
                render_path,
                mode=mode,
                ocr_language=ocr_language,
                dpi=int(dpi),
                remote=remote,
            )
        else:
            page_model = _analyze_raster_page(
                render_path, page_number, page_dir, source_mode="raster_image",
                mode=mode, ocr_language=ocr_language, dpi=int(dpi), remote=remote,
            )
        _write_json_atomic(model_path, page_model)
        receipt = {
            "schema_version": SCHEMA_VERSION,
            "state": "complete",
            "source_sha256": source_sha,
            "run_key": run_key,
            "page": page_number,
            "source_mode": page_model["source_mode"],
            "objects_path": str(model_path),
            "render_path": str(render_path),
            "metrics": page_model.get("metrics", {}),
            "completed_at": _utc_now(),
            "resumed": False,
        }
        _write_json_atomic(receipt_path, receipt)
        page_models.append(page_model)
        page_receipts.append(receipt)

    output.parent.mkdir(parents=True, exist_ok=True)
    build_target = output.with_name(f".{output.stem}.{os.getpid()}.tmp.pptx")
    if build_target.exists():
        build_target.unlink()
    assembly = _build_presentation(page_models, build_target)
    os.replace(build_target, output)
    output_sha = sha256_file(output)

    rendered_output = _render_pptx_for_comparison(output, qa_dir, int(dpi))
    comparisons = []
    for page_number, candidate in zip(selected_pages, rendered_output):
        reference = renders_dir / f"page-{page_number:04d}.png"
        comparisons.append({"page": page_number, **_comparison_metrics(reference, candidate)})

    try:
        final_visual_qa = run_visual_qa(
            output,
            pages="all",
            profile="slides",
            deterministic_only=True,
            receipt_root=run_dir / "visual-qa",
            resume=resume,
            dpi=int(dpi),
        )
    except Exception as exc:
        final_visual_qa = {
            "status": "partial",
            "verdict": "not_checked",
            "error": f"{type(exc).__name__}: {str(exc)[:500]}",
        }

    fallback_objects = sum(
        int(page_model.get("metrics", {}).get("fallback_object_count", 0))
        for page_model in page_models
    )
    coverage = [
        float(page_model.get("metrics", {}).get("editable_foreground_coverage", 0.0))
        for page_model in page_models
    ]
    status = "complete" if fallback_objects == 0 else "partial"
    summary = {
        "schema_version": SCHEMA_VERSION,
        "state": "complete",
        "status": status,
        "source": str(source),
        "source_sha256": source_sha,
        "output": str(output),
        "output_sha256": output_sha,
        "run_key": run_key,
        "run_dir": str(run_dir),
        "summary_path": str(summary_path),
        "selected_pages": selected_pages,
        "page_receipts": [str(pages_dir / f"page-{number:04d}" / "receipt.json") for number in selected_pages],
        "source_modes": [page_model["source_mode"] for page_model in page_models],
        "assembly": assembly,
        "fallback_objects": fallback_objects,
        "editable_foreground_coverage_mean": float(sum(coverage) / max(len(coverage), 1)),
        "fully_editable": fallback_objects == 0,
        "visual_comparison": comparisons,
        "visual_qa": final_visual_qa,
        "remote_allowed": remote_allowed,
        "remote_configured": remote_configured,
        "resumed_pages": sum(bool(receipt.get("resumed")) for receipt in page_receipts),
        "input_unchanged": sha256_file(source) == source_sha,
        "completed_at": _utc_now(),
        "resumed": False,
    }
    _write_json_atomic(summary_path, summary)
    _write_json_atomic(final_receipt_path, summary)
    return summary


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Convert an image or PDF into a resumable, editable PowerPoint reconstruction"
    )
    parser.add_argument("path")
    parser.add_argument("--output")
    parser.add_argument("--pages", default="all")
    parser.add_argument("--mode", choices=["faithful", "balanced", "editable"], default="balanced")
    parser.add_argument("--profile", choices=["auto", "figure", "paper", "slides", "document"], default="auto")
    parser.add_argument("--receipt-root")
    parser.add_argument("--ocr-language", default="auto")
    parser.add_argument("--dpi", type=int, default=180)
    parser.add_argument("--timeout", type=int, default=120)
    parser.add_argument("--model")
    parser.add_argument("--fallback-model")
    parser.add_argument("--no-resume", action="store_true")
    parser.add_argument("--overwrite", action="store_true")
    parser.add_argument("--allow-remote", action="store_true")
    parser.add_argument("--deterministic-only", action="store_true")
    return parser


def main() -> int:
    arguments = build_parser().parse_args()
    try:
        summary = run_visual_to_editable_ppt(
            arguments.path,
            output_path=arguments.output,
            pages=arguments.pages,
            mode=arguments.mode,
            profile=arguments.profile,
            receipt_root=arguments.receipt_root,
            ocr_language=arguments.ocr_language,
            dpi=arguments.dpi,
            timeout=arguments.timeout,
            model=arguments.model,
            fallback_model=arguments.fallback_model,
            resume=not arguments.no_resume,
            overwrite=arguments.overwrite,
            allow_remote=arguments.allow_remote,
            deterministic_only=arguments.deterministic_only,
        )
    except VisualReconstructionError as exc:
        print(json.dumps({"status": "failed", "error": str(exc)}, ensure_ascii=False))
        return 2
    print(json.dumps(summary, ensure_ascii=False, separators=(",", ":")))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
